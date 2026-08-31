from copy import deepcopy
from pptx import Presentation

SRC = "fab-mes-平台介绍.pptx"
prs = Presentation(SRC)

# --- idempotent: remove any existing NPI slide before re-adding ---
def slide_has_text(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return True
    return False

sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
to_remove = []
for idx, sl in enumerate(prs.slides):
    if slide_has_text(sl, "NPI THREAD"):
        to_remove.append(ids[idx])
for sldId in to_remove:
    prs.part.drop_rel(sldId.rId)
    sldIdLst.remove(sldId)
print("removed existing NPI slides:", len(to_remove))

# --- duplicate slide 11 (OTD digital thread) ---
src = list(prs.slides)[10]
new_slide = prs.slides.add_slide(src.slide_layout)
# clear any placeholder shapes the layout injected
for shp in list(new_slide.shapes):
    sp = shp._element
    sp.getparent().remove(sp)
# copy all shapes from source
for shp in src.shapes:
    new_slide.shapes._spTree.append(deepcopy(shp._element))

# --- helpers ---
def set_para_text(para, text):
    if para.runs:
        para.runs[0].text = text
        for r in para.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        r = para.add_run()
        r.text = text

def set_frame_lines(tf, lines):
    paras = tf.paragraphs
    for i, ln in enumerate(lines):
        if i < len(paras):
            set_para_text(paras[i], ln)
        else:
            p = tf.add_paragraph()
            set_para_text(p, ln)
    # drop any surplus paragraphs left over from the source shape
    for p in paras[len(lines):]:
        p._p.getparent().remove(p._p)

shapes = list(new_slide.shapes)

# title
for sh in shapes:
    if sh.has_text_frame and "DIGITAL THREAD" in sh.text_frame.text:
        set_frame_lines(sh.text_frame, ["NPI THREAD · 设计到流片"])
# subtitle
for sh in shapes:
    if sh.has_text_frame and "一条数字主线" in sh.text_frame.text:
        set_frame_lines(sh.text_frame, ["一条 NPI 主线，焊通设计 → 流片全生命周期"])
# accent caption (footer line1)
for sh in shapes:
    if sh.has_text_frame and "七大环节全部订阅" in sh.text_frame.text:
        set_frame_lines(sh.text_frame, ["↑ NPI 与 OTD 共享同一条 MES WS 事件总线（:8124）：design→route 自动派生，工程批/流片批沿同一产线推进。"])
# footer line2
for sh in shapes:
    if sh.has_text_frame and ("APC / VM / APS" in sh.text_frame.text or "APC / VM" in sh.text_frame.text):
        set_frame_lines(sh.text_frame, ["npi-ops.html 统一管理台：设计档案 / 光罩 / 工程批 / 流片批，一键投放，状态实时可见。"])

# 7 cards in order
cards = [sh for sh in shapes if sh.has_text_frame and sh.text_frame.text.startswith(("①","②","③","④","⑤","⑥","⑦"))]
card_texts = [
    ["① 设计档案", "MES · DB", "GDSII·PDK 主数据"],
    ["② 光罩", "photomasks", "层数→passes 派生"],
    ["③ 工程批", "engineering WO", "NPI 一键投放"],
    ["④ 流片", "tapeout WO", "首跑+qualification 重入"],
    ["⑤ MES 主轴量测", "MES · :8124", "与 OTD 同引擎"],
    ["⑥ 判异 SPC·FDC", "五大引擎", "实时质量判异"],
    ["⑦ 转量产 ramp", "并入 OTD", "设计→交付闭环"],
]
assert len(cards) == 7, f"cards={len(cards)}"
for sh, lines in zip(cards, card_texts):
    set_frame_lines(sh.text_frame, lines)

# --- reorder: place new slide right after original slide 11 (index 10 -> insert at 11) ---
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
new_id = ids[-1]
sldIdLst.remove(new_id)
sldIdLst.insert(11, new_id)

prs.save(SRC)
print("saved; total slides now:", len(prs.slides))
