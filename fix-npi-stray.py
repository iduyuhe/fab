from pptx import Presentation

SRC = "fab-mes-平台介绍.pptx"
OUT = "fab-mes-平台介绍.clean.pptx"
prs = Presentation(SRC)

def find_npi(slides):
    for s in slides:
        for sh in s.shapes:
            if sh.has_text_frame and "NPI THREAD" in sh.text_frame.text:
                return s
    return None

npi = find_npi(prs.slides)
assert npi is not None, "NPI slide not found"
fixed = 0
for sh in npi.shapes:
    if sh.has_text_frame and "npi-ops.html" in sh.text_frame.text:
        paras = sh.text_frame.paragraphs
        # keep only first paragraph; drop surplus (stray "投料 → ..." tail)
        for p in paras[1:]:
            p._p.getparent().remove(p._p)
            fixed += 1
print("footer paragraphs removed:", fixed)

prs.save(OUT)
print("saved ->", OUT)
